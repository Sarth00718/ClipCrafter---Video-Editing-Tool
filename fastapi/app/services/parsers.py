"""
Document parsers for PDF, DOCX, PPTX, and TXT files.

Each parser returns a list of dictionaries with keys:
    - text: str          — extracted text content
    - page_number: int   — 1-indexed page / slide number (None for TXT)
    - section_heading: str | None — detected heading (best-effort)
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Protocol

from app.core.logger import get_logger
from app.utils.helpers import clean_text

logger = get_logger(__name__)


# ────────────────────────────── Base Protocol ──────────────────────


class DocumentParser(Protocol):
    """Interface every parser must satisfy."""

    def parse(self, file_path: Path) -> list[dict]:
        """Parse a file and return a list of page-level dicts."""
        ...


# ────────────────────────────── PDF Parser ─────────────────────────


class PDFParser:
    """Extract text from PDF files using PyMuPDF (fitz)."""

    def parse(self, file_path: Path) -> list[dict]:
        import fitz  # PyMuPDF

        logger.info("Parsing PDF: %s", file_path.name)
        pages: list[dict] = []

        try:
            doc = fitz.open(str(file_path))
            for page_num, page in enumerate(doc, start=1):
                raw = page.get_text("text")
                text = clean_text(raw)
                if text:
                    heading = self._detect_heading(text)
                    pages.append({
                        "text": text,
                        "page_number": page_num,
                        "section_heading": heading,
                    })
            doc.close()
        except Exception as exc:
            logger.error("Failed to parse PDF %s: %s", file_path.name, exc)
            raise

        logger.info("Extracted %d pages from PDF", len(pages))
        return pages

    @staticmethod
    def _detect_heading(text: str) -> str | None:
        """Heuristic: first short line that looks like a heading."""
        for line in text.splitlines()[:5]:
            stripped = line.strip()
            if 3 < len(stripped) < 120 and not stripped.endswith("."):
                return stripped
        return None


# ────────────────────────────── DOCX Parser ────────────────────────


class DOCXParser:
    """Extract text from DOCX files using python-docx."""

    def parse(self, file_path: Path) -> list[dict]:
        from docx import Document as DocxDocument

        logger.info("Parsing DOCX: %s", file_path.name)
        pages: list[dict] = []

        try:
            doc = DocxDocument(str(file_path))
            current_heading: str | None = None
            section_text_lines: list[str] = []
            section_index = 1

            for para in doc.paragraphs:
                # Detect heading styles
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    # Flush previous section
                    if section_text_lines:
                        pages.append({
                            "text": clean_text("\n".join(section_text_lines)),
                            "page_number": section_index,
                            "section_heading": current_heading,
                        })
                        section_index += 1
                        section_text_lines = []
                    current_heading = para.text.strip()
                else:
                    if para.text.strip():
                        section_text_lines.append(para.text)

            # Flush last section
            if section_text_lines:
                pages.append({
                    "text": clean_text("\n".join(section_text_lines)),
                    "page_number": section_index,
                    "section_heading": current_heading,
                })
        except Exception as exc:
            logger.error("Failed to parse DOCX %s: %s", file_path.name, exc)
            raise

        logger.info("Extracted %d sections from DOCX", len(pages))
        return pages


# ────────────────────────────── PPTX Parser ────────────────────────


class PPTXParser:
    """Extract text from PowerPoint files using python-pptx."""

    def parse(self, file_path: Path) -> list[dict]:
        from pptx import Presentation

        logger.info("Parsing PPTX: %s", file_path.name)
        slides: list[dict] = []

        try:
            prs = Presentation(str(file_path))
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                title: str | None = None

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = paragraph.text.strip()
                            if line:
                                texts.append(line)

                    # Capture slide title
                    if shape.has_text_frame and shape.shape_id == 1:
                        title = shape.text_frame.text.strip() or None

                combined = clean_text("\n".join(texts))
                if combined:
                    slides.append({
                        "text": combined,
                        "page_number": None,
                        "slide_number": slide_num,
                        "section_heading": title,
                    })
        except Exception as exc:
            logger.error("Failed to parse PPTX %s: %s", file_path.name, exc)
            raise

        logger.info("Extracted %d slides from PPTX", len(slides))
        return slides


# ────────────────────────────── TXT Parser ─────────────────────────


class TXTParser:
    """Read plain-text files."""

    def parse(self, file_path: Path) -> list[dict]:
        logger.info("Parsing TXT: %s", file_path.name)
        try:
            raw = file_path.read_text(encoding="utf-8", errors="replace")
            text = clean_text(raw)
        except Exception as exc:
            logger.error("Failed to read TXT %s: %s", file_path.name, exc)
            raise

        if not text:
            return []

        return [{
            "text": text,
            "page_number": 1,
            "section_heading": None,
        }]


# ────────────────────────────── CSV Parser ─────────────────────────


class CSVParser:
    """Extract text from CSV files using pandas."""

    def parse(self, file_path: Path) -> list[dict]:
        import pandas as pd

        logger.info("Parsing CSV: %s", file_path.name)
        try:
            df = pd.read_csv(str(file_path))
            
            # Convert DataFrame to readable text format
            text_parts = []
            
            # Add column headers
            text_parts.append("Columns: " + ", ".join(df.columns))
            
            # Add row data
            for idx, row in df.iterrows():
                row_text = " | ".join([f"{col}: {val}" for col, val in row.items()])
                text_parts.append(row_text)
            
            combined_text = clean_text("\n".join(text_parts))
            
            return [{
                "text": combined_text,
                "page_number": 1,
                "section_heading": "CSV Data",
            }]
        except Exception as exc:
            logger.error("Failed to parse CSV %s: %s", file_path.name, exc)
            raise


# ────────────────────────────── XLSX Parser ────────────────────────


class XLSXParser:
    """Extract text from Excel files using openpyxl."""

    def parse(self, file_path: Path) -> list[dict]:
        import pandas as pd

        logger.info("Parsing XLSX: %s", file_path.name)
        pages: list[dict] = []
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(str(file_path))
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names, start=1):
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert sheet to text
                text_parts = []
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append("Columns: " + ", ".join(df.columns))
                
                for idx, row in df.iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items()])
                    text_parts.append(row_text)
                
                combined_text = clean_text("\n".join(text_parts))
                
                if combined_text:
                    pages.append({
                        "text": combined_text,
                        "page_number": sheet_num,
                        "section_heading": sheet_name,
                    })
            
            logger.info("Extracted %d sheets from XLSX", len(pages))
            return pages
            
        except Exception as exc:
            logger.error("Failed to parse XLSX %s: %s", file_path.name, exc)
            raise


# ────────────────────────────── MD Parser ──────────────────────────


class MDParser:
    """Extract text from Markdown files."""

    def parse(self, file_path: Path) -> list[dict]:
        logger.info("Parsing MD: %s", file_path.name)
        pages: list[dict] = []
        
        try:
            raw = file_path.read_text(encoding="utf-8", errors="replace")
            
            # Split by headers (# Header)
            sections = re.split(r'\n(?=#{1,6}\s)', raw)
            
            for idx, section in enumerate(sections, start=1):
                section = section.strip()
                if not section:
                    continue
                
                # Extract heading
                heading_match = re.match(r'^#{1,6}\s+(.+?)$', section, re.MULTILINE)
                heading = heading_match.group(1) if heading_match else None
                
                # Clean text
                text = clean_text(section)
                
                if text:
                    pages.append({
                        "text": text,
                        "page_number": idx,
                        "section_heading": heading,
                    })
            
            # If no sections found, treat as single page
            if not pages:
                text = clean_text(raw)
                if text:
                    pages.append({
                        "text": text,
                        "page_number": 1,
                        "section_heading": None,
                    })
            
            logger.info("Extracted %d sections from MD", len(pages))
            return pages
            
        except Exception as exc:
            logger.error("Failed to parse MD %s: %s", file_path.name, exc)
            raise


# ────────────────────────────── Factory ────────────────────────────


_PARSER_MAP: dict[str, type] = {
    "pdf": PDFParser,
    "docx": DOCXParser,
    "pptx": PPTXParser,
    "txt": TXTParser,
    "csv": CSVParser,
    "xlsx": XLSXParser,
    "md": MDParser,
}


class ParserFactory:
    """Return the appropriate parser for a given file extension."""

    @staticmethod
    def get_parser(extension: str) -> DocumentParser:
        ext = extension.lower().lstrip(".")
        parser_cls = _PARSER_MAP.get(ext)
        if parser_cls is None:
            raise ValueError(f"Unsupported file type: .{ext}")
        return parser_cls()
